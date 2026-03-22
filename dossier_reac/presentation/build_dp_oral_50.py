#!/usr/bin/env python3
"""Build a 50-slide HappyRow oral presentation from the repaired template."""

from pathlib import Path
from textwrap import dedent

from PIL import Image
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

SCRIPT_DIR = Path(__file__).resolve().parent
REPO_ROOT = SCRIPT_DIR.parent.parent
TEMPLATE = SCRIPT_DIR / "ppt" / "HappyRow-Template  -  Repaired.pptx"
OUTPUT = SCRIPT_DIR / "ppt" / "HappyRow_DP_50_Slides_Jimmy_Ni.pptx"

FRONT_ASSETS = Path("/Users/j.ni/IdeaProjects/happyrow-front/design-figma")
MERISE_ASSETS = REPO_ROOT / "dossier_reac" / "full_project" / "merise"

IMG_LOGO = FRONT_ASSETS / "logo.png"
IMG_START = FRONT_ASSETS / "start_page.png"
IMG_HOME = FRONT_ASSETS / "home_page.png"
IMG_CREATE_EVENT = FRONT_ASSETS / "create new event form.png"
IMG_EVENT = FRONT_ASSETS / "event_page.png"
IMG_MCD = MERISE_ASSETS / "mcd.png"
IMG_MPD = MERISE_ASSETS / "mpd.png"

NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"
NS_R = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"


def note(text: str) -> str:
    return " ".join(dedent(text).strip().split())


def cover_slide(title, subtitle_lines, notes, image=None, image_box=None):
    return {
        "kind": "cover",
        "layout": "COVER",
        "title": title,
        "subtitle_lines": subtitle_lines,
        "notes": notes,
        "image": image,
        "image_box": image_box,
    }


def main_slide(title, bullets, notes, layout="MAIN 1", font_size=None, title_size=24):
    return {
        "kind": "main",
        "layout": layout,
        "title": title,
        "bullets": bullets,
        "notes": notes,
        "font_size": font_size,
        "title_size": title_size,
    }


def split_text_slide(
    title,
    left_bullets,
    right_bullets,
    notes,
    layout="SPLIT 1",
    left_font_size=None,
    right_font_size=None,
    title_size=24,
):
    return {
        "kind": "split_text",
        "layout": layout,
        "title": title,
        "left_bullets": left_bullets,
        "right_bullets": right_bullets,
        "notes": notes,
        "left_font_size": left_font_size,
        "right_font_size": right_font_size,
        "title_size": title_size,
    }


def split_image_slide(
    title,
    bullets,
    image,
    notes,
    layout="SPLIT 1",
    image_side="right",
    font_size=None,
    title_size=24,
):
    return {
        "kind": "split_image",
        "layout": layout,
        "title": title,
        "bullets": bullets,
        "image": image,
        "notes": notes,
        "image_side": image_side,
        "font_size": font_size,
        "title_size": title_size,
    }


SLIDES = [
    cover_slide(
        "HappyRow",
        [
            "Plateforme collaborative de gestion d'evenements",
            "Presentation du dossier de projet CDA",
            "Jimmy Ni - Session du 24/03/2026",
        ],
        notes=note(
            """
            Bonjour, je vais vous presenter HappyRow, le projet full-stack que j'ai realise
            dans le cadre du titre professionnel Concepteur Developpeur d'Applications.
            L'objectif de cette soutenance est de montrer comment j'ai couvert l'ensemble
            des competences attendues, depuis l'analyse du besoin jusqu'au deploiement en
            production. Je vais parcourir le contexte, les choix de conception, les
            realisations marquantes, puis la securite, les tests et enfin le bilan du projet.
            """
        ),
        image=IMG_LOGO,
        image_box=(13.6, 0.7, 3.2, 1.4),
    ),
    main_slide(
        "Plan de la presentation",
        [
            "Le besoin metier et le cadre du projet",
            "Le pilotage, la qualite et le deploiement",
            "La conception fonctionnelle et technique",
            "Les realisations les plus significatives",
            "La securite, les tests, la veille et le bilan",
        ],
        notes=note(
            """
            Je vais suivre une progression assez naturelle. Je commencerai par le besoin
            auquel HappyRow repond et par le cadre de formation dans lequel le projet a ete
            mene. J'enchainerai ensuite sur la gestion de projet et les mecanismes de
            qualite, avant de presenter la conception fonctionnelle puis les specifications
            techniques. Je terminerai par les extraits de realisation les plus representatifs,
            la securite, la strategie de test, la veille et enfin la conclusion.
            """
        ),
        layout="MAIN 1",
        font_size=19,
    ),
    split_image_slide(
        "HappyRow en quelques mots",
        [
            "Une application web collaborative pour organiser un evenement a plusieurs",
            "Une plateforme qui centralise participants, ressources et contributions",
            "Deux repositories : un back-end Kotlin/Ktor et un front-end React/TypeScript",
            "Un projet pense pour demontrer les competences du referentiel CDA",
        ],
        image=IMG_HOME,
        notes=note(
            """
            HappyRow part d'un probleme tres concret : lorsqu'on organise un evenement a
            plusieurs, on se retrouve vite avec des messages disperses, des listes
            informelles et peu de visibilite sur ce que chacun apporte. J'ai donc concu une
            application web collaborative qui centralise l'organisation d'un evenement, la
            gestion des participants, les ressources attendues et les contributions
            individuelles. Le projet est decoupe en deux repositories complementaires : une
            API REST Kotlin/Ktor et une interface utilisateur React en TypeScript.
            """
        ),
        layout="SPLIT 1",
        font_size=16,
    ),
    main_slide(
        "Les competences CDA mobilisees",
        [
            "CCP 1 : environnement, interfaces, composants metier et gestion de projet",
            "CCP 2 : analyse, architecture hexagonale, base de donnees et acces SQL",
            "CCP 3 : tests, deploiement, CI/CD et mise en production",
            "Le projet a ete construit pour couvrir les 11 competences du referentiel",
        ],
        notes=note(
            """
            Dans le dossier, j'ai fait le choix de relier explicitement HappyRow aux trois
            CCP du titre CDA. Le premier bloc porte sur le developpement d'une application
            securisee, avec l'environnement de travail, l'interface utilisateur et les
            composants metier. Le second se concentre sur la conception en couches, donc
            l'analyse, l'architecture hexagonale et la modelisation de la base. Enfin, le
            troisieme bloc couvre les tests, le deploiement, la documentation et la mise en
            production dans une logique DevOps.
            """
        ),
        layout="MAIN 2",
        font_size=18,
    ),
    main_slide(
        "Le besoin auquel repond le projet",
        [
            "L'organisation d'un anniversaire, d'un diner ou d'une fete est souvent eparpillee",
            "Les messages et tableurs improvises donnent peu de visibilite",
            "Il manque un outil dedie pour savoir qui vient et qui apporte quoi",
            "HappyRow apporte structure, partage et suivi collaboratif",
        ],
        notes=note(
            """
            Le point de depart, c'est un besoin d'usage tres simple. Quand un groupe
            prepare un anniversaire, un diner ou une fete, les informations circulent dans
            plusieurs canaux et deviennent vite difficiles a suivre. On ne sait plus tres
            bien qui participe, quelles ressources sont encore necessaires ou si plusieurs
            personnes se sont positionnees sur la meme chose. HappyRow a donc ete pense
            comme un outil specialise qui centralise cette coordination et la rend lisible.
            """
        ),
        layout="MAIN 3",
        font_size=18,
    ),
    main_slide(
        "Les objectifs fonctionnels",
        [
            "Permettre a un utilisateur authentifie de creer et gerer des evenements",
            "Inviter des participants et suivre leur statut",
            "Definir les ressources necessaires avec categorie et quantite suggeree",
            "Permettre a chacun de declarer sa contribution",
            "Garantir la coherence des donnees en cas d'acces concurrents",
        ],
        notes=note(
            """
            A partir de ce besoin, j'ai formalise plusieurs objectifs fonctionnels. Le
            premier est de permettre a un utilisateur authentifie de creer, consulter,
            modifier et supprimer ses evenements. Le deuxieme est d'inviter des
            participants et de suivre leur statut. J'ai ensuite ajoute la gestion des
            ressources necessaires a l'evenement, puis la declaration des contributions.
            Enfin, j'ai voulu traiter un vrai sujet technique, a savoir la coherence des
            donnees lorsque plusieurs personnes contribuent en meme temps.
            """
        ),
        layout="MAIN 1",
        font_size=18,
    ),
    split_text_slide(
        "Le perimetre fonctionnel",
        [
            "Evenements : creer, consulter, modifier et supprimer",
            "Participants : ajouter, lister et changer le statut",
        ],
        [
            "Ressources : creer, consulter et suivre les quantites",
            "Contributions : ajouter, reduire ou supprimer avec verrou optimiste",
        ],
        notes=note(
            """
            Le perimetre fonctionnel s'articule autour de quatre domaines. D'abord les
            evenements, avec un CRUD complet reserve aux utilisateurs authentifies. Ensuite
            les participants, qui peuvent etre ajoutes par email et dont le statut peut
            evoluer. Le troisieme domaine concerne les ressources necessaires a
            l'evenement, par exemple de la nourriture, des boissons ou du materiel. Enfin,
            le quatrieme domaine est celui des contributions, avec la possibilite d'ajouter,
            de reduire ou de supprimer une contribution tout en protegeant les acces
            concurrents.
            """
        ),
        layout="SPLIT 2",
        left_font_size=16,
        right_font_size=16,
    ),
    main_slide(
        "Contraintes techniques et livrables",
        [
            "Kotlin sur JVM, Ktor, PostgreSQL et authentification JWT via Supabase",
            "Docker, Render, architecture hexagonale et qualite via Detekt et Spotless",
            "Une API REST de 15 endpoints couvrant 4 domaines metier",
            "Une base PostgreSQL complete, de la documentation et des tests",
            "Un front-end web consommant l'API dans un repository dedie",
        ],
        notes=note(
            """
            Le cahier des charges impose aussi un cadre technique assez clair. Le back-end
            doit etre realise en Kotlin avec Ktor, s'appuyer sur PostgreSQL et deleguer
            l'authentification a Supabase via des JWT. L'architecture attendue est
            hexagonale, le deploiement doit passer par Docker et Render, et la qualite doit
            etre controlee avec de l'analyse statique et du formatage. En sortie, je dois
            livrer une API REST, une base relationnelle coherente, de la documentation,
            des tests et une application front-end qui consomme cette API.
            """
        ),
        layout="MAIN 2",
        font_size=17,
    ),
    main_slide(
        "Contexte de formation",
        [
            "Projet realise dans la formation CDA DevOps chez Simplon",
            "Promotion de 20 apprenants, organisee a distance",
            "Projet personnel full-stack mene pendant les temps dedies au projet",
            "Objectif pedagogique : demontrer les competences du referentiel CDA",
        ],
        notes=note(
            """
            HappyRow n'est pas un projet commande par une entreprise cliente, mais un
            projet personnel realise dans le cadre de ma formation Concepteur Developpeur
            d'Applications DevOps chez Simplon. La formation se deroule a distance et
            repose sur des projets individuels et collectifs. Dans ce contexte, HappyRow
            m'a servi de support principal pour demontrer ma capacite a analyser un besoin,
            construire une application complete et l'amener jusqu'au deploiement en
            production.
            """
        ),
        layout="MAIN 3",
        font_size=18,
    ),
    main_slide(
        "Environnement humain et responsabilites",
        [
            "Je porte la conception, le developpement, les tests et le deploiement",
            "J'ai une autonomie complete sur les choix techniques et architecturaux",
            "Les formateurs interviennent en conseil, revue de projet et validation",
            "Les echanges avec la promotion servent surtout au partage de bonnes pratiques",
        ],
        notes=note(
            """
            Sur ce projet, mon role est central puisqu'il s'agit d'un travail individuel.
            J'ai assume la conception de l'architecture, le developpement du back-end et
            du front-end, la modelisation de la base, les tests et le deploiement. Les
            formateurs ont joue un role d'accompagnement et de validation, mais sans
            intervenir directement dans le code. Cette organisation m'a oblige a prendre
            les decisions techniques, a les argumenter et a les faire converger vers un
            resultat coherent.
            """
        ),
        layout="MAIN 1",
        font_size=18,
    ),
    main_slide(
        "Une gestion de projet en mode Kanban",
        [
            "Visualiser le flux To Do, In Progress et Done",
            "Limiter le travail en cours pour rester focalise",
            "Livrer de maniere continue et incrementale",
            "Traiter chaque ticket comme une fonctionnalite, un correctif ou une tache technique",
        ],
        notes=note(
            """
            Pour piloter HappyRow, j'ai choisi une approche Kanban. Ce choix est adapte a
            un projet mene seul, parce qu'il apporte de la visibilite sans imposer un
            formalisme trop lourd. Le tableau me permet de suivre l'avancement, de limiter
            le nombre de sujets ouverts en meme temps et de faire avancer le projet par
            increments fonctionnels. Chaque ticket correspond soit a une fonctionnalite,
            soit a un bug, soit a une tache technique, ce qui rend le pilotage tres lisible.
            """
        ),
        layout="MAIN 2",
        font_size=18,
    ),
    main_slide(
        "Decoupage en cinq phases",
        [
            "Fondations : architecture, Gradle, Docker, PostgreSQL et CI/CD",
            "Evenements : CRUD complet et authentification JWT Supabase",
            "Participants : auto-ajout du createur et gestion des statuts",
            "Ressources et contributions : mecanisme de verrou optimiste",
            "Qualite et deploiement : tests, analyse statique et mise en ligne",
        ],
        notes=note(
            """
            J'ai egalement structure le projet en grandes phases. La premiere consiste a
            poser les fondations techniques avec l'architecture, le build multi-modules, la
            base et les premiers pipelines. La seconde porte sur le domaine evenement. La
            troisieme ajoute la gestion des participants. La quatrieme introduit les
            ressources et surtout les contributions avec verrou optimiste. Enfin, la
            cinquieme phase consolide l'ensemble avec les tests, l'analyse statique et le
            deploiement en production.
            """
        ),
        layout="MAIN 3",
        font_size=17,
    ),
    main_slide(
        "Outils de suivi et workflow Git",
        [
            "GitHub Issues pour suivre les taches, bugs et ameliorations",
            "GitHub Projects pour visualiser l'avancement dans le tableau Kanban",
            "Git pour les branches de fonctionnalite et les merges par pull request",
            "GitHub Actions pour obtenir un feedback automatise a chaque push",
        ],
        notes=note(
            """
            Les outils de suivi s'articulent tres logiquement avec la methode choisie.
            GitHub Issues me sert de point d'entree pour les besoins et les corrections,
            tandis que GitHub Projects me donne la vision Kanban. Cote versioning, je
            travaille avec une branche principale stable et des branches de fonctionnalite
            isolees. Enfin, GitHub Actions joue un role de filet de securite, puisqu'il
            me retourne automatiquement les resultats de qualite, de tests et de build.
            """
        ),
        layout="MAIN 1",
        font_size=18,
    ),
    main_slide(
        "Pipeline CI/CD front-end",
        [
            "Job Test : lockfile-lint, npm ci, ESLint, Vitest et build",
            "Job Build and Push Docker pour produire l'image front-end",
            "Job Approval pour valider la mise en production",
            "Job Deploy sur Vercel en mode production",
            "Workflow Security Audit quotidien sur la supply chain",
        ],
        notes=note(
            """
            Le front-end dispose de son propre pipeline GitHub Actions. La premiere etape
            verifie la chaine de dependances et execute le lint, les tests unitaires puis le
            build. La deuxieme construit et pousse l'image Docker. La troisieme ajoute une
            approbation manuelle avant la mise en production, ce qui introduit une
            validation humaine. Enfin, le deploiement s'effectue sur Vercel. En parallele,
            j'ai aussi mis en place une veille automatisee sur la supply chain avec un
            workflow quotidien de securite.
            """
        ),
        layout="MAIN 2",
        font_size=17,
    ),
    main_slide(
        "Pipeline CI/CD back-end et deploiement",
        [
            "Sequence Detekt, tests unitaires, deploiement Render puis notification",
            "Chaque push sur main declenche automatiquement le pipeline",
            "Render rebuild l'image Docker multi-stage et redeploie le service",
            "Les secrets et variables d'environnement sont geres hors du code",
            "Le conteneur d'execution reste minimal et non-root",
        ],
        notes=note(
            """
            Le pipeline back-end est plus lineaire. A chaque push sur la branche main,
            GitHub Actions execute d'abord Detekt, puis les tests unitaires. Si ces etapes
            passent, le pipeline notifie Render qui reconstruit l'image Docker multi-stage
            et redeploie le service. Ce flux automatise me permet d'enchainer qualite et
            mise en production sans intervention manuelle. J'ai aussi isole les secrets et
            les variables d'environnement dans les services prevus a cet effet, pour garder
            le code source propre et limiter l'exposition.
            """
        ),
        layout="MAIN 3",
        font_size=17,
    ),
    main_slide(
        "Objectifs qualite et environnement technique",
        [
            "Analyse statique avec Detekt et formatage homogone via Spotless",
            "Tests automatises et pipelines bloquants avant deploiement",
            "Documentation du code par nommage explicite et organisation claire",
            "Stack principale : Kotlin 2.2, Ktor 3.2.2, React 19, TypeScript strict et Docker",
        ],
        notes=note(
            """
            Au-dela des fonctionnalites, j'ai voulu installer des exigences de qualite
            assez explicites. L'analyse statique et le formatage sont integres au flux de
            travail pour eviter la derive du code. Les pipelines servent de garde-fou avant
            chaque mise en ligne. J'ai aussi travaille la lisibilite du code par le nommage
            et par une organisation en couches. Enfin, l'environnement technique repose
            sur des versions recentes et coherentes entre back-end, front-end et
            infrastructure.
            """
        ),
        layout="MAIN 1",
        font_size=17,
    ),
    main_slide(
        "Architecture hexagonale : le choix directeur",
        [
            "Deux modules Gradle independants : domain et infrastructure",
            "Un domaine sans dependance technique ni framework",
            "Des adapters REST, SQL et JWT regroupes dans l'infrastructure",
            "Une injection de dependances assuree par Koin",
        ],
        notes=note(
            """
            L'architecture hexagonale est le fil conducteur du projet. J'ai separe le
            coeur metier dans un module domain, sans dependance vers Ktor, Exposed ou toute
            autre bibliotheque technique. Tout ce qui touche au transport HTTP, a la
            persistence SQL ou a l'authentification est place dans le module infrastructure.
            Cette separation rend les dependances explicites et facilite la testabilite.
            Koin me sert ensuite a assembler les briques au moment de l'execution.
            """
        ),
        layout="MAIN 2",
        font_size=18,
    ),
    main_slide(
        "Principe de dependance et bounded contexts",
        [
            "Le domaine definit les ports, l'infrastructure les implemente",
            "La logique metier reste testable sans la base ni le framework web",
            "Le projet est decoupe en Event, Resource, Participant et Contribution",
            "Changer d'ORM ou de framework ne remet pas en cause le coeur metier",
        ],
        notes=note(
            """
            Le point cle de cette architecture, c'est l'inversion de dependance. Les use
            cases du domaine ne connaissent que des interfaces de repository, donc des
            ports. Les implementations concretes sont apportees par l'infrastructure. Ce
            principe permet de tester la logique metier isolee, mais aussi de raisonner par
            concepts metier plutot que par couches techniques. C'est pour cela que j'ai
            organise le projet autour de quatre bounded contexts : evenement, participant,
            ressource et contribution.
            """
        ),
        layout="MAIN 3",
        font_size=17,
    ),
    main_slide(
        "Flux d'une requete dans l'application",
        [
            "L'endpoint recoit et valide la requete HTTP",
            "Les DTO convertissent l'entree vers le domaine",
            "Le use case appelle les repositories via les ports",
            "L'infrastructure persiste en base et reconstruit la reponse",
        ],
        notes=note(
            """
            Si je prends maintenant une requete type, son parcours illustre bien la
            separation des responsabilites. L'endpoint Ktor recoit la requete, la
            deserialise et controle sa validite. Les DTO assurent le passage entre monde
            HTTP et monde domaine. Le use case applique ensuite la regle metier et dialogue
            avec les repositories abstraits. Enfin, l'infrastructure realise les operations
            SQL et renvoie une representation de sortie, ce qui permet de conserver un flux
            clair et testable de bout en bout.
            """
        ),
        layout="MAIN 1",
        font_size=18,
    ),
    split_image_slide(
        "Maquettes, design tokens et navigation",
        [
            "Des maquettes Figma realisees avant le developpement",
            "Une charte graphique basee sur teal, navy, coral et Comic Neue",
            "Une navigation SPA avec React Router DOM",
            "Un rendu conditionnel pilote par l'etat d'authentification",
        ],
        image=IMG_START,
        notes=note(
            """
            Avant de coder, j'ai formalise l'interface avec des maquettes Figma. Cela m'a
            permis de clarifier les ecrans attendus, les interactions et l'identite visuelle
            du projet. J'ai ensuite traduit cette charte en design tokens CSS pour garder
            une coherence globale. Cote front-end, l'application fonctionne comme une SPA
            avec React Router DOM, et le rendu des ecrans depend directement de l'etat
            d'authentification de l'utilisateur.
            """
        ),
        layout="SPLIT 1",
        font_size=16,
    ),
    split_image_slide(
        "Les ecrans principaux de l'application",
        [
            "WelcomeView pour l'arrivee des utilisateurs non authentifies",
            "LoginModal et RegisterModal pour l'entree dans l'application",
            "HomePage comme tableau de bord des evenements",
            "UserProfilePage pour le profil et la deconnexion",
        ],
        image=IMG_CREATE_EVENT,
        notes=note(
            """
            Dans le rapport, j'ai detaille les principaux ecrans qui structurent le parcours
            utilisateur. On commence par une page d'accueil publique, puis les modales de
            connexion et d'inscription permettent l'acces a l'application. Une fois
            authentifie, l'utilisateur retrouve la HomePage qui joue le role de dashboard
            des evenements. Le profil et la deconnexion sont isoles dans une vue dediee,
            ce qui garde la navigation simple et coherente pour un usage quotidien.
            """
        ),
        layout="SPLIT 2",
        font_size=16,
    ),
    split_image_slide(
        "L'ecran de detail d'evenement et la contribution",
        [
            "Une vue complete avec categories de ressources",
            "Des controles plus et moins pour preparer un delta de contribution",
            "Un bouton Validate qui n'apparait qu'au moment utile",
            "L'ajout de ressource inline et la suppression d'evenement sur la meme vue",
        ],
        image=IMG_EVENT,
        notes=note(
            """
            Cette vue est la plus representative de l'application, parce qu'elle concentre
            a la fois l'information metier et l'interaction utilisateur. On y retrouve les
            ressources organisees par categorie, avec la quantite courante, la quantite
            suggeree et les controles de contribution. L'utilisateur peut preparer un delta
            puis le valider explicitement. J'ai aussi choisi d'autoriser certaines actions
            directement dans cette vue, comme l'ajout de ressource, afin de limiter les
            allers-retours et garder une experience fluide.
            """
        ),
        layout="SPLIT 1",
        font_size=16,
    ),
    split_image_slide(
        "Modele conceptuel de donnees",
        [
            "Quatre entites principales : Event, Participant, Resource et Contribution",
            "Un evenement heberge plusieurs participants et plusieurs ressources",
            "La contribution relie un participant a une ressource",
            "Le modele traduit un besoin collaboratif avec des relations simples",
        ],
        image=IMG_MCD,
        notes=note(
            """
            La modelisation des donnees repose sur quatre entites qui correspondent
            directement aux concepts metier. Un evenement regroupe des participants et des
            ressources, tandis que la contribution sert de lien entre une personne et une
            ressource donnee. Cette modelisation est volontairement compacte, mais elle
            couvre l'essentiel du besoin fonctionnel. Elle constitue aussi une bonne base
            pour mettre en oeuvre des contraintes d'integrite et un traitement fiable des
            contributions concurrentes.
            """
        ),
        layout="SPLIT 2",
        font_size=16,
    ),
    split_image_slide(
        "Du modele logique au modele physique",
        [
            "Le MLD materialise les cles etrangeres entre les tables",
            "Le MPD precise les types PostgreSQL, les enums et les valeurs par defaut",
            "Le schema cible reste centre sur 4 tables principales",
            "Un champ version est ajoute sur Resource pour le verrou optimiste",
        ],
        image=IMG_MPD,
        notes=note(
            """
            Une fois le modele conceptuel pose, je l'ai traduit en modele logique puis en
            modele physique PostgreSQL. Les associations deviennent des cles etrangeres,
            les types sont precises, les contraintes sont explicitees et certains comportements
            sont portes par des valeurs par defaut. J'ai aussi introduit un champ version
            sur la table Resource, car il est indispensable pour porter le mecanisme de
            verrou optimiste. Cette etape montre bien le passage d'une vision metier a une
            implementation relationnelle concrete.
            """
        ),
        layout="SPLIT 1",
        font_size=16,
    ),
    main_slide(
        "Contraintes d'integrite et index",
        [
            "Cles primaires UUID generees automatiquement",
            "Cles etrangeres pour relier evenement, participant, ressource et contribution",
            "Indexes uniques composites pour eviter les doublons metier",
            "CHECK sur les quantites et valeurs par defaut sur status et version",
            "Indexes dedies aux recherches frequentes pour les performances",
        ],
        notes=note(
            """
            La base ne se contente pas de stocker les donnees, elle participe aussi a la
            coherence metier. Les cles primaires garantissent l'identification, les cles
            etrangeres maintiennent les liens entre entites, et les indexes uniques
            composites evitent certains doublons metier comme deux contributions identiques
            sur une meme ressource. J'ai ajoute des contraintes CHECK et des valeurs par
            defaut pour renforcer encore cette robustesse. Enfin, les indexes de recherche
            ameliorent les acces frequents sans alourdir la logique applicative.
            """
        ),
        layout="MAIN 2",
        font_size=17,
    ),
    main_slide(
        "Le verrou optimiste au coeur des contributions",
        [
            "La ressource porte une version technique qui evolue a chaque mise a jour",
            "L'UPDATE verifie la version attendue avant de modifier la quantite",
            "En cas de conflit, la requete echoue et renvoie un 409 Conflict",
            "Ce choix est prefere a un verrou pessimiste dans ce contexte",
        ],
        notes=note(
            """
            Le mecanisme technique le plus interessant du projet est le verrou optimiste.
            L'idee est de ne pas bloquer la ressource en base, mais de verifier qu'entre la
            lecture et l'ecriture personne n'a deja modifie sa quantite. Concretement, la
            table Resource porte un champ version. La mise a jour ne reussit que si cette
            version correspond toujours a celle attendue. Sinon, l'operation est refusee,
            une exception metier est levee et l'API retourne un 409 Conflict au client.
            """
        ),
        layout="MAIN 3",
        font_size=18,
    ),
    main_slide(
        "Les 12 cas d'utilisation du projet",
        [
            "Creer, consulter, modifier et supprimer un evenement",
            "Ajouter et suivre les participants",
            "Creer et consulter les ressources d'un evenement",
            "Ajouter, reduire ou supprimer une contribution",
            "Auto-ajouter le createur comme participant confirme",
        ],
        notes=note(
            """
            Le diagramme de cas d'utilisation synthese le perimetre fonctionnel autour de
            douze cas principaux. On y retrouve naturellement la gestion des evenements,
            des participants, des ressources et des contributions. J'y ai egalement fait
            apparaitre des regles metier importantes, notamment l'auto-ajout du createur
            comme participant confirme, ou encore la suppression d'ensemble qui accompagne
            la suppression d'un evenement. Ce diagramme me sert de pont entre l'expression
            du besoin et la decomposition du code en use cases.
            """
        ),
        layout="MAIN 1",
        font_size=17,
    ),
    main_slide(
        "Sequence de creation d'un evenement",
        [
            "Reception de la requete HTTP de creation",
            "Deserialisation, validation puis appel du use case",
            "Insertion de l'evenement dans la base",
            "Auto-ajout du createur comme participant confirme",
            "Retour d'un DTO de reponse vers le client",
        ],
        notes=note(
            """
            Cette sequence montre le parcours complet d'une creation d'evenement. L'API
            recoit d'abord la requete et verifie qu'elle est exploitable. Le use case
            orchestre ensuite la creation metier, puis applique immediatement la regle qui
            ajoute le createur comme participant confirme. Cela garantit qu'un evenement ne
            se retrouve jamais sans organisateur rattache. Enfin, la reponse est retransformee
            en DTO pour le client, ce qui maintient la separation entre monde HTTP et monde
            metier.
            """
        ),
        layout="MAIN 2",
        font_size=17,
    ),
    main_slide(
        "Sequence d'ajout d'une contribution",
        [
            "Authentification de l'utilisateur via son JWT",
            "Recherche ou creation du participant rattache a l'evenement",
            "Verification d'une contribution existante",
            "Mise a jour atomique de la ressource et de sa version",
            "Gestion du conflit si deux utilisateurs agissent en meme temps",
        ],
        notes=note(
            """
            La sequence d'ajout d'une contribution est la plus riche du projet, car elle
            traverse plusieurs couches et combine regles metier, securite et concurrence.
            L'utilisateur est d'abord authentifie, puis l'application retrouve ou cree le
            participant correspondant a son identite. Elle verifie ensuite s'il existe deja
            une contribution sur cette ressource, puis met a jour la quantite et la version
            de la ressource de maniere atomique. Si un autre utilisateur a modifie la
            ressource entre-temps, le conflit est detecte proprement.
            """
        ),
        layout="MAIN 3",
        font_size=17,
    ),
    main_slide(
        "Stack back-end",
        [
            "Kotlin 2.2 et JVM 21 pour la robustesse et la null-safety",
            "Ktor 3.2.2 pour le serveur HTTP asynchrone",
            "Exposed, PostgreSQL et HikariCP pour la persistence",
            "Koin pour l'injection et Arrow pour les erreurs fonctionnelles",
            "Auth0 JWT et Supabase pour la verification des tokens",
        ],
        notes=note(
            """
            La stack back-end a ete choisie pour concilier modernite, lisibilite et
            robustesse. Kotlin m'apporte un typage fort, une bonne expressivite et la
            gestion de la null-safety. Ktor repond bien au besoin d'une API legere et
            asynchrone. Exposed et PostgreSQL couvrent la persistence relationnelle, tandis
            que HikariCP gere efficacement les connexions. J'ai ajoute Koin pour l'injection
            de dependances et Arrow pour modeliser les erreurs metier de facon fonctionnelle.
            """
        ),
        layout="MAIN 1",
        font_size=17,
    ),
    main_slide(
        "Stack front-end",
        [
            "React 19 et TypeScript strict pour l'interface",
            "Vite pour le build et React Router DOM pour la navigation",
            "Supabase JS SDK pour l'authentification et la session",
            "PWA avec Workbox et service worker",
            "CSS natif et design tokens, sans framework CSS lourd",
        ],
        notes=note(
            """
            Cote front-end, j'ai privilegie une stack moderne mais volontairement sobre.
            React 19 et TypeScript me donnent un bon compromis entre productivite,
            fiabilite et lisibilite. Vite accelere le developpement et le build, tandis que
            React Router DOM structure la navigation de la SPA. L'authentification repose
            sur le SDK Supabase, et la PWA ajoute l'installabilite ainsi qu'un mode offline.
            Enfin, j'ai choisi de rester sur du CSS natif avec design tokens pour garder la
            maitrise de la charte sans introduire de framework visuel supplementaire.
            """
        ),
        layout="MAIN 2",
        font_size=17,
    ),
    main_slide(
        "Infrastructure commune et hebergement",
        [
            "Docker multi-stage pour le back-end et le front-end",
            "GitHub Actions pour les pipelines d'integration et de deploiement",
            "Back-end heberge sur Render, a Francfort",
            "Front-end heberge sur Vercel",
            "IntelliJ IDEA comme IDE principal de developpement",
        ],
        notes=note(
            """
            Les deux applications partagent une logique d'infrastructure commune. J'ai
            conteneurise le back-end et le front-end avec Docker en multi-stage afin de
            separer clairement les images de build et de runtime. GitHub Actions assure le
            lien entre qualite, tests et deploiement. Le back-end est heberge sur Render,
            tandis que le front-end est deploye sur Vercel. Ce couple d'outils me permet
            d'obtenir une mise en ligne simple, reproductible et adaptee a un projet de
            cette taille.
            """
        ),
        layout="MAIN 3",
        font_size=17,
    ),
    main_slide(
        "Les patterns de conception appliques",
        [
            "Use Case : une classe dediee pour chaque operation metier",
            "Repository : interface metier et implementation SQL",
            "DTO : separation explicite entre domaine et API",
            "Either d'Arrow pour composer les erreurs sans exceptions",
            "Value class Creator pour renforcer le typage",
        ],
        notes=note(
            """
            Plusieurs patterns structurent le projet. Les use cases portent chacun une
            responsabilite metier precise. Les repositories permettent de separer les
            besoins du domaine et la facon d'acceder aux donnees. Les DTO servent de zone
            tampon entre l'API et le coeur metier. Avec Arrow Either, j'ai choisi une
            gestion d'erreurs explicite et composable plutot qu'un usage massif des
            exceptions. Enfin, certaines value classes me permettent d'exprimer davantage
            d'intention metier tout en gardant une implementation legere.
            """
        ),
        layout="MAIN 1",
        font_size=17,
    ),
    main_slide(
        "Les principes SOLID dans HappyRow",
        [
            "Single Responsibility : un use case, un endpoint, un repository par role",
            "Open/Closed : de nouveaux contextes peuvent etre ajoutes sans casser l'existant",
            "Liskov et Interface Segregation via des contrats de repository cibles",
            "Dependency Inversion : le domaine depend d'abstractions, pas d'implementations",
        ],
        notes=note(
            """
            Au-dela des patterns, j'ai essaye de faire vivre les principes SOLID dans la
            structure du projet. Le principe de responsabilite unique se voit dans le
            decoupage des classes. L'ouverture a l'extension sans modification excessive est
            rendue possible par les bounded contexts et les abstractions. Les interfaces de
            repository restent specifiques a leur besoin pour eviter des contrats trop
            larges. Et surtout, l'inversion de dependance garantit que le domaine reste le
            point d'ancrage de l'architecture.
            """
        ),
        layout="MAIN 2",
        font_size=17,
    ),
    main_slide(
        "API REST : conventions et endpoints",
        [
            "Base path : /event/configuration/api/v1",
            "Ressources au pluriel et identifiants en path parameters",
            "Endpoints publics : /, /info et /health",
            "15 endpoints couvrant evenements, participants, ressources et contributions",
            "Reponses JSON en snake_case",
        ],
        notes=note(
            """
            L'API REST suit des conventions simples et stables. J'utilise un base path
            versionne, des ressources nommees au pluriel et des identifiants dans l'URL.
            Trois routes restent publiques pour les verifications de base, alors que les
            routes metier necessitent une authentification. Le perimetre total represente
            quinze endpoints qui couvrent les quatre domaines fonctionnels du projet. Enfin,
            les reponses sont produites en JSON avec une convention snake_case, ce qui
            apporte de la coherence cote API.
            """
        ),
        layout="MAIN 3",
        font_size=17,
    ),
    main_slide(
        "Gestion des erreurs HTTP",
        [
            "201 pour une creation, 200 pour un succes et 204 pour une suppression",
            "400 pour une requete invalide ou mal formee",
            "401 pour un token manquant ou invalide",
            "403 pour une action interdite, 404 pour une ressource absente",
            "409 pour un conflit de nom ou de concurrence",
        ],
        notes=note(
            """
            J'ai aussi formalise une grille de lecture claire des erreurs HTTP. Les codes
            de succes distinguent creation, operation standard et suppression. Les erreurs
            client couvrent les corps de requete invalides, l'absence ou l'invalidite du
            token, les actions non autorisees et les ressources introuvables. Le 409 a une
            place importante dans HappyRow, car il sert a signaler a la fois certains
            conflits metier et surtout l'echec du verrou optimiste en cas d'acces
            concurrent.
            """
        ),
        layout="MAIN 1",
        font_size=17,
    ),
    main_slide(
        "Les choix d'eco-conception",
        [
            "Un perimetre fonctionnel volontairement sobre et centre sur le besoin reel",
            "Peu de dependances de production cote front-end",
            "Des images Docker optimisees grace au multi-stage",
            "Une PWA et du cache pour reduire les requetes repetees",
            "Une JVM et un pool de connexions dimensionnes au besoin",
        ],
        notes=note(
            """
            Meme si ce n'est pas un projet d'eco-conception au sens strict, j'ai voulu
            integrer quelques principes de sobriete numerique. J'ai limite le perimetre
            fonctionnel aux usages utiles, evite les dependances lourdes cote front-end et
            optimise les images Docker avec du multi-stage. Le mode PWA et le cache des
            assets reduisent aussi certains appels repetes. Enfin, le dimensionnement de la
            JVM et du pool de connexions reste proportionne a la taille du service, ce qui
            correspond bien a l'esprit du projet.
            """
        ),
        layout="MAIN 2",
        font_size=17,
    ),
    main_slide(
        "Securite dans les specifications techniques",
        [
            "JWT Supabase verifie a chaque requete",
            "Validation des entrees cote serveur",
            "CORS en liste blanche",
            "Secrets uniquement en variables d'environnement",
            "HTTPS et SSL base de donnees en production",
        ],
        notes=note(
            """
            La securite est prise en compte des la phase de specification technique. Toutes
            les routes metier reposent sur une verification du JWT Supabase. Les corps et
            parametres de requete sont valides cote serveur, ce qui evite de faire
            confiance au client. La configuration CORS reste restrictive avec une liste
            blanche d'origines. Les secrets ne sont jamais stockes dans le code source, et
            la mise en production impose des canaux securises aussi bien pour HTTP que pour
            la connexion a la base.
            """
        ),
        layout="MAIN 3",
        font_size=17,
    ),
    main_slide(
        "Pourquoi ne pas avoir integre de NoSQL",
        [
            "Les donnees sont fortement relationnelles et riches en contraintes d'integrite",
            "Le projet a besoin de transactions ACID pour les operations critiques",
            "Une base NoSQL aurait ajoute une complexite d'infrastructure peu justifiee",
            "PostgreSQL couvre deja certains besoins semi-structures via JSONB et arrays",
            "Une ouverture reste possible pour une future evolution du projet",
        ],
        notes=note(
            """
            Le referentiel CDA mentionne les acces SQL et NoSQL, mais dans HappyRow j'ai
            choisi de rester sur PostgreSQL. Ce choix vient d'abord de la nature des
            donnees, qui sont tres relationnelles et fortement contraintes. Ensuite, le
            verrou optimiste et plusieurs operations atomiques s'appuient naturellement sur
            les garanties transactionnelles d'un SGBD relationnel. Introduire une base NoSQL
            n'aurait pas apporte de gain immediat au regard du perimetre. En revanche, j'ai
            garde cette perspective ouverte pour des usages futurs comme le cache ou les logs.
            """
        ),
        layout="MAIN 1",
        font_size=17,
    ),
    split_image_slide(
        "Realisation UI : la fonctionnalite de contribution",
        [
            "Chaque ressource affiche quantite courante et quantite suggeree",
            "L'utilisateur selectionne un delta avec les boutons plus et moins",
            "La validation est explicite avant l'envoi a l'API",
            "L'interface privilegie la visibilite et la simplicite d'usage",
        ],
        image=IMG_EVENT,
        notes=note(
            """
            Dans les realisations, j'ai choisi de mettre en avant la fonctionnalite de
            contribution car elle concentre les enjeux metier du projet. L'utilisateur voit
            immediatement ce qui est deja prevu, ce qu'il reste a apporter et peut ajuster
            sa participation avec un delta simple. Le fait de rendre la validation explicite
            permet de garder le controle sur l'action envoyee au serveur. Cette interface
            fait donc le lien entre une interaction accessible et un traitement metier plus
            technique en arriere-plan.
            """
        ),
        layout="SPLIT 2",
        font_size=16,
    ),
    main_slide(
        "Realisation back-end : endpoint de creation d'evenement",
        [
            "L'endpoint recoit la requete et securise l'entree HTTP",
            "Either.catch transforme les erreurs de deserialisation",
            "flatMap n'enchaine la logique que si la requete est valide",
            "toDomain() et toDto() isolent l'API du coeur metier",
        ],
        notes=note(
            """
            Cet extrait illustre bien la facon dont j'ai organise les endpoints Ktor.
            L'objectif n'est pas de mettre de la logique metier dans la couche HTTP, mais de
            lui faire jouer un role de facade propre. Either.catch me permet d'intercepter
            les erreurs de deserialisation sans casser le flux. Ensuite, flatMap garantit
            que le use case n'est appele que lorsque les entrees sont valides. Enfin, la
            conversion entre DTO et domaine maintient une frontiere nette entre contrat API
            et logique applicative.
            """
        ),
        layout="MAIN 2",
        font_size=17,
    ),
    main_slide(
        "Realisation back-end : gestion du 409 Conflict",
        [
            "La cause racine est remontee pour distinguer conflit et erreur technique",
            "Un 409 est retourne si la ressource a change entre lecture et ecriture",
            "Le message reste clair pour le client sans exposer les details internes",
            "La gestion d'erreur reste coherente du repository jusqu'a l'endpoint",
        ],
        notes=note(
            """
            J'ai voulu traiter le conflit de concurrence comme un cas metier a part entiere,
            pas comme une simple erreur technique. Pour cela, la cause racine est analysee
            afin de distinguer un echec de verrou optimiste d'un probleme purement
            technique. L'endpoint peut alors renvoyer un 409 Conflict avec un message
            intelligible pour le client. Cette remontee d'information est importante parce
            qu'elle garde du sens metier tout en evitant de divulguer des details internes
            sur la base ou l'implementation.
            """
        ),
        layout="MAIN 3",
        font_size=17,
    ),
    main_slide(
        "Realisation front-end : composant ResourceItem",
        [
            "Le composant choisit entre ajout, mise a jour ou suppression",
            "L'etat local selectedQuantity represente un delta",
            "Le bouton de validation n'apparait que lorsqu'une action est en cours",
            "L'interaction reste lisible avant tout appel reseau",
        ],
        notes=note(
            """
            Cote front-end, ResourceItem est un bon exemple d'intelligence locale dans
            l'interface. Le composant ne se contente pas d'afficher des donnees : il porte
            aussi la logique qui determine s'il faut creer, modifier ou supprimer une
            contribution en fonction du resultat attendu. Le state local selectedQuantity
            est traite comme un delta, ce qui rend l'interaction tres intuitive. Et comme
            le bouton de validation n'apparait qu'au bon moment, l'utilisateur percoit
            clairement quand une action va vraiment etre envoyee au serveur.
            """
        ),
        layout="MAIN 1",
        font_size=17,
    ),
    main_slide(
        "Realisation front-end : repository HTTP et JWT",
        [
            "Les appels API reposent sur fetch natif, sans dependance tierce",
            "Le token est obtenu via une callback getToken injectee",
            "Le mapping snake_case vers camelCase est centralise",
            "Le decouplage entre UI et HTTP est preserve",
        ],
        notes=note(
            """
            J'ai egalement mis en avant le repository HTTP front-end, car il montre que
            j'ai cherche a reproduire une separation en couches comparable a celle du
            back-end. Les appels API passent par fetch natif, ce qui limite les
            dependances. Le token JWT n'est pas lu de facon implicite dans le repository :
            il est injecte via une callback, ce qui rend la dependance explicite. Enfin,
            les transformations entre snake_case et camelCase sont centralisees, ce qui
            simplifie les composants d'interface et les use cases du front.
            """
        ),
        layout="MAIN 2",
        font_size=17,
    ),
    main_slide(
        "Realisation metier : CreateEventUseCase",
        [
            "Le use case orchestre creation d'evenement et auto-ajout du createur",
            "Il depend uniquement des ports EventRepository et ParticipantRepository",
            "flatMap garantit la coherence de l'enchainement",
            "mapLeft enveloppe les erreurs avec un contexte metier",
        ],
        notes=note(
            """
            CreateEventUseCase est un bon resume de mon approche du domaine. Il orchestre
            deux operations qui doivent rester coherentes : la creation de l'evenement puis
            l'auto-ajout du createur comme participant confirme. Comme le use case ne
            depend que de ports, il reste independant des details techniques. L'utilisation
            de flatMap me permet d'exprimer clairement la dependance entre les etapes,
            tandis que mapLeft enrichit les erreurs pour garder un vrai contexte metier.
            """
        ),
        layout="MAIN 3",
        font_size=17,
    ),
    main_slide(
        "Realisation metier : entites et modele de domaine",
        [
            "Des data class Kotlin pour l'immutabilite et la lisibilite",
            "Aucune annotation de framework dans le domaine",
            "Creator encapsule une chaine dans une value class",
            "La propriete version de Resource porte la semantique du verrou optimiste",
        ],
        notes=note(
            """
            Le modele de domaine a ete travaille pour rester le plus pur possible.
            L'absence d'annotations framework evite de polluer les concepts metier avec des
            contraintes de persistence ou de transport. Les data class Kotlin rendent les
            objets faciles a manipuler tout en restant lisibles. La value class Creator
            renforce l'intention metier a peu de frais. Et surtout, la propriete version sur
            Resource permet de faire exister le verrou optimiste directement dans le modele,
            pas seulement dans la couche SQL.
            """
        ),
        layout="MAIN 1",
        font_size=17,
    ),
    main_slide(
        "Acces aux donnees : repository SQL et table Contribution",
        [
            "Le repository gere ajout, reduction et suppression de contribution",
            "La ressource associee est mise a jour dans la meme logique metier",
            "Un uniqueIndex composite empeche les doublons de contribution",
            "Une contrainte CHECK garantit des quantites strictement positives",
            "Cles etrangeres et indexes soutiennent integrite et performances",
        ],
        notes=note(
            """
            Dans la couche d'acces aux donnees, le repository SQL de contribution est l'un
            des composants les plus riches. Il doit gerer plusieurs variantes de mise a jour
            tout en restant coherent avec l'etat de la ressource associee. La table
            Contribution complete cette logique avec des contraintes fortes : unicite du
            couple participant-ressource, quantites positives et liens referentiels explicites.
            Cet ensemble montre bien comment j'ai combine regles metier, contraintes base de
            donnees et preoccupations de performance dans une meme implementation.
            """
        ),
        layout="MAIN 2",
        font_size=16,
    ),
    main_slide(
        "Securite de l'application",
        [
            "Authentification JWT cote API et gestion de session cote client",
            "Autorisation basee sur l'identite extraite du token",
            "Protection contre l'injection SQL via Exposed et requetes parametrees",
            "Gestion d'erreurs securisee, CORS restrictif et secrets hors du code",
            "Defense en profondeur entre endpoint, domaine, base et infrastructure",
        ],
        notes=note(
            """
            Le chapitre securite synthese les protections mises en place a plusieurs
            niveaux. L'authentification repose sur les JWT Supabase verifies a chaque
            requete. Les regles d'autorisation se basent sur l'identite extraite du token,
            pas sur des informations envoyees librement par le client. Les entrees sont
            validees, les requetes SQL sont parametrees, le CORS reste controle et les
            secrets sont externalises. Au final, j'ai cherche a mettre en place une defense
            en profondeur qui combine couche HTTP, domaine, base et infrastructure.
            """
        ),
        layout="MAIN 3",
        font_size=17,
    ),
    main_slide(
        "Tests, veille securite et apprentissages",
        [
            "Tests back-end unitaires et d'integration, plus analyse statique continue",
            "Cote front-end : 46 tests unitaires et 6 tests E2E selon le dossier",
            "Jeu d'essai centre sur la contribution avec verrou optimiste",
            "Veille via OWASP, CVE Database, GitHub Advisories, documentations et ANSSI",
            "Recherche specifique sur les tests de concurrence avec Testcontainers",
        ],
        notes=note(
            """
            La qualite ne s'arrete pas au code produit. J'ai defini une strategie de tests
            en pyramide, avec des tests unitaires, des tests d'integration, de l'analyse
            statique et des tests E2E cote front-end. Le jeu d'essai le plus representatif
            porte sur l'ajout de contribution avec verrou optimiste, parce qu'il mobilise
            authentification, logique metier, persistence et gestion d'erreur. En parallele,
            j'ai mene une veille securite continue sur les sources OWASP, CVE, GitHub
            Advisories, Ktor, Supabase et ANSSI. Cette veille m'a notamment aide a traiter
            la question delicate des tests de concurrence avec Testcontainers.
            """
        ),
        layout="MAIN 1",
        font_size=16,
    ),
    main_slide(
        "Bilan, perspectives et conclusion",
        [
            "Une application full-stack fonctionnelle et deployee en production",
            "Des satisfactions fortes autour de l'architecture, d'Either, du verrou optimiste et du CI/CD",
            "Des difficultes techniques sur la concurrence et la suppression en cascade",
            "Des perspectives : mobile, notifications, NoSQL, monitoring et tests d'integration en CI",
            "Un projet qui synthese les competences acquises pendant la formation",
        ],
        notes=note(
            """
            Pour conclure, HappyRow m'a permis de mener un projet complet, de l'analyse du
            besoin jusqu'a la mise en production. J'en retiens plusieurs points de
            satisfaction : l'architecture hexagonale, l'usage d'Arrow Either, le mecanisme
            de verrou optimiste et l'automatisation apportee par les pipelines CI/CD. Les
            principales difficultes ont porte sur la gestion des acces concurrents et sur la
            suppression en cascade des donnees liees. Enfin, j'ai deja identifie plusieurs
            pistes d'evolution comme le mobile, les notifications, le monitoring et
            l'integration plus poussee des tests d'integration dans la CI.
            """
        ),
        layout="MAIN 2",
        font_size=16,
    ),
]


def layout_by_name(prs: Presentation, name: str):
    for layout in prs.slide_layouts:
        if layout.name == name:
            return layout
    raise KeyError(f"Unknown layout: {name}")


def placeholder_by_idx(slide, idx: int):
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == idx:
            return shape
    raise KeyError(f"Missing placeholder idx {idx} on slide {slide.slide_id}")


def default_font_size(bullets, split=False):
    if split:
        if len(bullets) >= 4:
            return 16
        return 17
    if len(bullets) >= 5:
        return 16
    if len(bullets) == 4:
        return 17
    return 18


def write_title(shape, text: str, size=24, align=None):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    paragraph = tf.paragraphs[0]
    paragraph.text = text
    paragraph.alignment = align or PP_ALIGN.LEFT
    paragraph.font.size = Pt(size)
    paragraph.font.bold = True


def write_lines(shape, lines, size=18, align=PP_ALIGN.CENTER):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    for index, line in enumerate(lines):
        paragraph = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        paragraph.text = line
        paragraph.alignment = align
        paragraph.font.size = Pt(size if index == 0 else max(size - 2, 12))
        paragraph.space_before = Pt(8 if index else 0)


def write_bullets(shape, bullets, size):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    for index, bullet in enumerate(bullets):
        paragraph = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        paragraph.text = bullet
        paragraph.level = 0
        paragraph.font.size = Pt(size)
        paragraph.space_before = Pt(0)
        paragraph.space_after = Pt(8 if len(bullets) <= 4 else 6)
        paragraph.line_spacing = 1.05


def add_notes(slide, text: str):
    slide.notes_slide.notes_text_frame.text = text


def add_picture_contained(slide, path: Path, box):
    if not path.exists():
        print(f"Warning: missing image {path}")
        return

    left, top, width, height = box
    left = Inches(left)
    top = Inches(top)
    width = Inches(width)
    height = Inches(height)

    with Image.open(path) as image:
        img_width, img_height = image.size

    img_ratio = img_width / img_height
    box_ratio = width / height
    if img_ratio > box_ratio:
        final_width = width
        final_height = int(width / img_ratio)
        final_left = left
        final_top = top + int((height - final_height) / 2)
    else:
        final_height = height
        final_width = int(height * img_ratio)
        final_top = top
        final_left = left + int((width - final_width) / 2)

    slide.shapes.add_picture(str(path), final_left, final_top, final_width, final_height)


def add_picture_on_placeholder(slide, placeholder_idx, path: Path, padding=0.12):
    placeholder = placeholder_by_idx(slide, placeholder_idx)
    box = (
        placeholder.left / 914400 + padding,
        placeholder.top / 914400 + padding,
        placeholder.width / 914400 - (2 * padding),
        placeholder.height / 914400 - (2 * padding),
    )
    add_picture_contained(slide, path, box)


def render_slide(prs: Presentation, data: dict):
    slide = prs.slides.add_slide(layout_by_name(prs, data["layout"]))

    if data["kind"] == "cover":
        write_title(placeholder_by_idx(slide, 0), data["title"], size=28, align=PP_ALIGN.CENTER)
        write_lines(placeholder_by_idx(slide, 10), data["subtitle_lines"], size=18)
        if data.get("image") and data.get("image_box"):
            add_picture_contained(slide, data["image"], data["image_box"])

    elif data["kind"] == "main":
        write_title(placeholder_by_idx(slide, 0), data["title"], size=data["title_size"])
        font_size = data["font_size"] or default_font_size(data["bullets"], split=False)
        write_bullets(placeholder_by_idx(slide, 10), data["bullets"], font_size)

    elif data["kind"] == "split_text":
        write_title(placeholder_by_idx(slide, 0), data["title"], size=data["title_size"])
        left_size = data["left_font_size"] or default_font_size(data["left_bullets"], split=True)
        right_size = data["right_font_size"] or default_font_size(data["right_bullets"], split=True)
        write_bullets(placeholder_by_idx(slide, 10), data["left_bullets"], left_size)
        write_bullets(placeholder_by_idx(slide, 11), data["right_bullets"], right_size)

    elif data["kind"] == "split_image":
        write_title(placeholder_by_idx(slide, 0), data["title"], size=data["title_size"])
        font_size = data["font_size"] or default_font_size(data["bullets"], split=True)
        if data["image_side"] == "right":
            write_bullets(placeholder_by_idx(slide, 10), data["bullets"], font_size)
            add_picture_on_placeholder(slide, 11, data["image"])
        else:
            add_picture_on_placeholder(slide, 10, data["image"])
            write_bullets(placeholder_by_idx(slide, 11), data["bullets"], font_size)

    add_notes(slide, data["notes"])


def remove_original_template_slides(prs: Presentation, count: int):
    slide_id_list = prs.element.find(f"{{{NS_P}}}sldIdLst")
    originals = list(slide_id_list)[:count]
    for slide_id in originals:
        rel_id = slide_id.get(f"{{{NS_R}}}id")
        slide_id_list.remove(slide_id)
        if rel_id:
            try:
                prs.part.drop_rel(rel_id)
            except Exception:
                pass


def build():
    if len(SLIDES) != 50:
        raise ValueError(f"Expected 50 slides, got {len(SLIDES)}")

    prs = Presentation(str(TEMPLATE))
    original_count = len(prs.slides)

    for slide_data in SLIDES:
        render_slide(prs, slide_data)

    remove_original_template_slides(prs, original_count)
    prs.save(str(OUTPUT))

    notes_count = sum(
        1
        for slide in prs.slides
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text.strip()
    )
    print(f"Created {OUTPUT} with {len(prs.slides)} slides and {notes_count} notes.")


if __name__ == "__main__":
    build()
