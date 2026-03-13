#!/usr/bin/env python3
"""Post-process the pandoc PPTX to remove overflow slides and clean up."""

from pptx import Presentation
from pptx.util import Pt, Emu, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from copy import deepcopy
from lxml import etree
import sys
import os

NAVY = RGBColor(0x3D, 0x5A, 0x6C)
TEAL = RGBColor(0x5F, 0xBD, 0xB4)
CORAL = RGBColor(0xE6, 0xA1, 0x9A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x2D, 0x3A, 0x3A)

def get_slide_title(slide):
    """Extract title text from slide."""
    for shape in slide.shapes:
        if shape.has_text_frame:
            if hasattr(shape, 'placeholder_format') and shape.placeholder_format is not None:
                if shape.placeholder_format.idx == 0:
                    return shape.text_frame.text.strip()
    for shape in slide.shapes:
        if shape.has_text_frame and shape.shape_id == 2:
            return shape.text_frame.text.strip()
    return ""


def is_overflow_slide(slide):
    """Detect if a slide is an overflow (no title, auto-created by pandoc)."""
    title = get_slide_title(slide)
    if title:
        return False
    has_content = False
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text_frame.text.strip():
            has_content = True
            break
        if shape.shape_type is not None and shape.shape_type == 13:
            has_content = True
            break
    return has_content


def remove_slide(prs, slide):
    """Remove a slide from the presentation by manipulating the XML."""
    rId = None
    for rel in prs.part.rels.values():
        if rel.target_part == slide.part:
            rId = rel.rId
            break
    if rId:
        ns_p = '{http://schemas.openxmlformats.org/presentationml/2006/main}'
        ns_r = '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}'
        sldIdLst = prs.element.find(f'{ns_p}sldIdLst')
        for sldId in list(sldIdLst):
            if sldId.get(f'{ns_r}id') == rId:
                sldIdLst.remove(sldId)
                break
        prs.part.drop_rel(rId)


def style_title_slide(slide):
    """Apply HappyRow styling to title slide."""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        if hasattr(shape, 'placeholder_format') and shape.placeholder_format is not None:
            idx = shape.placeholder_format.idx
            if idx == 0:
                for para in shape.text_frame.paragraphs:
                    para.alignment = PP_ALIGN.CENTER
                    for run in para.runs:
                        run.font.size = Pt(36)
                        run.font.color.rgb = NAVY
                        run.font.bold = True
            elif idx == 1:
                for para in shape.text_frame.paragraphs:
                    para.alignment = PP_ALIGN.CENTER
                    for run in para.runs:
                        run.font.size = Pt(18)
                        run.font.color.rgb = TEAL


def style_content_slides(prs):
    """Apply HappyRow styling to content slides."""
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            if hasattr(shape, 'placeholder_format') and shape.placeholder_format is not None:
                idx = shape.placeholder_format.idx
                if idx == 0:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            run.font.color.rgb = NAVY
                            run.font.bold = True


def main():
    input_path = sys.argv[1] if len(sys.argv) > 1 else 'HappyRow.pptx'
    output_path = sys.argv[2] if len(sys.argv) > 2 else input_path

    prs = Presentation(input_path)
    total_before = len(prs.slides)

    overflow_slides = []
    for i, slide in enumerate(prs.slides):
        title = get_slide_title(slide)
        if not title and i > 0:
            overflow_slides.append(slide)

    print(f"Found {len(overflow_slides)} overflow slides to remove")
    for slide in overflow_slides:
        remove_slide(prs, slide)

    if prs.slides and len(prs.slides) > 0:
        style_title_slide(prs.slides[0])

    style_content_slides(prs)

    prs.save(output_path)

    total_after = len(prs.slides)
    print(f"Cleanup: {total_before} → {total_after} slides (removed {total_before - total_after} overflow slides)")
    print(f"Output: {output_path}")


if __name__ == '__main__':
    main()
