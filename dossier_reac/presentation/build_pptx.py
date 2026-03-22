#!/usr/bin/env python3
"""Build HappyRow.pptx from markdown slides using the custom template.

Usage: python3 build_pptx.py [--no-sections]
"""

import os
import re
import random
import sys
from copy import deepcopy
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

SCRIPT_DIR = Path(__file__).parent
SLIDES_DIR = SCRIPT_DIR / "slides"
TEMPLATE = SCRIPT_DIR / "happyrow_template.pptx"
NOTES_FILE = SCRIPT_DIR / "00_SCRIPT_COMPLET.md"
IMG_DIR = SCRIPT_DIR / "mermaid-images"
OUTPUT = SCRIPT_DIR / "HappyRow.pptx"

NAVY = RGBColor(0x3D, 0x5A, 0x6C)
TEAL = RGBColor(0x5F, 0xBD, 0xB4)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x2D, 0x3A, 0x3A)
GRAY = RGBColor(0x99, 0x99, 0x99)
CODE_BG = RGBColor(0xF5, 0xF5, 0xF5)

MAX_CODE_LINES = 22
NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"
NS_R = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"

T_TITLE = 0
T_SPLIT = [1, 2]
T_MAIN = [3, 4, 5]
T_CONCL = 6
T_SECTION = 7

random.seed(42)


# ---------------------------------------------------------------------------
# Markdown parsing
# ---------------------------------------------------------------------------

def parse_rich(text):
    """Split markdown bold/italic/code into (text, bold, italic, code) segments."""
    segs = []
    pat = re.compile(r"(\*\*(.+?)\*\*)|(\*(.+?)\*)|(`(.+?)`)")
    pos = 0
    for m in pat.finditer(text):
        if m.start() > pos:
            segs.append((text[pos : m.start()], False, False, False))
        if m.group(2):
            segs.append((m.group(2), True, False, False))
        elif m.group(4):
            segs.append((m.group(4), False, True, False))
        elif m.group(6):
            segs.append((m.group(6), False, False, True))
        pos = m.end()
    if pos < len(text):
        segs.append((text[pos:], False, False, False))
    return segs or [(text, False, False, False)]


def strip_md(text):
    t = re.sub(r"\*\*(.+?)\*\*", r"\1", text)
    t = re.sub(r"\*(.+?)\*", r"\1", t)
    t = re.sub(r"`(.+?)`", r"\1", t)
    return t.strip()


def parse_slides(slides_dir):
    """Return a list of {'type':'section'|'slide', ...} items."""
    items = []
    mermaid_idx = 0

    for md in sorted(Path(slides_dir).glob("*.md")):
        lines = md.read_text(encoding="utf-8").split("\n")
        cur = None
        in_code = in_mermaid = False
        code_lang = ""
        code_buf = []
        file_section_seen = False

        for raw in lines:
            s = raw.strip()
            if s.startswith("<!--") or s == "---":
                continue
            if s.startswith("> *[") or s.startswith(">*["):
                continue

            # --- h1: section header (first in file) or content heading ---
            mh = re.match(r"^# (.+)$", s)
            if mh and not s.startswith("## ") and not in_code and not in_mermaid:
                if not file_section_seen:
                    file_section_seen = True
                    if cur:
                        items.append(cur)
                        cur = None
                    items.append({"type": "section", "title": mh.group(1)})
                elif cur:
                    cur["content"].append({"type": "heading1", "text": mh.group(1)})
                continue

            # --- slide header ---
            ms = re.match(r"^## Slide (\d+)\s*[—–-]\s*(.+?)(?:\s*\(.+\))?$", s)
            if ms and not in_code and not in_mermaid:
                if cur:
                    items.append(cur)
                cur = {
                    "type": "slide",
                    "number": int(ms.group(1)),
                    "title": ms.group(2).strip(),
                    "content": [],
                    "has_code": False,
                    "has_image": False,
                }
                continue

            if cur is None:
                continue

            # --- mermaid ---
            if s == "```mermaid" and not in_code:
                in_mermaid = True
                mermaid_idx += 1
                continue
            if in_mermaid:
                if s == "```":
                    in_mermaid = False
                    p = str(IMG_DIR / f"diagram_{mermaid_idx}.png")
                    if os.path.exists(p):
                        cur["content"].append({"type": "image", "path": p})
                        cur["has_image"] = True
                continue

            # --- code block ---
            mc = re.match(r"^```(\w+)", s)
            if mc and not in_code:
                in_code = True
                code_lang = mc.group(1)
                code_buf = []
                continue
            if s == "```" and in_code:
                in_code = False
                cur["has_code"] = True
                cur["content"].append(
                    {"type": "code", "language": code_lang, "text": "\n".join(code_buf)}
                )
                continue
            if in_code:
                code_buf.append(raw)
                continue

            # --- plain code fence (no language) ---
            if s == "```":
                continue

            # --- table row ---
            if s.startswith("|") and s.endswith("|"):
                if re.match(r"^\|[\s\-:|]+\|$", s):
                    continue
                cells = [c.strip() for c in s.split("|")[1:-1]]
                if cur["content"] and cur["content"][-1]["type"] == "table":
                    cur["content"][-1]["rows"].append(cells)
                else:
                    cur["content"].append({"type": "table", "rows": [cells]})
                continue

            # --- sub-heading ---
            m3 = re.match(r"^### (.+)$", s)
            if m3:
                cur["content"].append({"type": "subheading", "text": m3.group(1)})
                continue

            # --- h2 inside slide (not a slide header) ---
            m2 = re.match(r"^## (.+)$", s)
            if m2:
                cur["content"].append({"type": "heading2", "text": m2.group(1)})
                continue

            # --- bullet ---
            if s.startswith("- "):
                cur["content"].append({"type": "bullet", "text": s[2:]})
                continue

            # --- numbered ---
            mn = re.match(r"^(\d+)\.\s+(.+)$", s)
            if mn:
                cur["content"].append(
                    {"type": "numbered", "num": mn.group(1), "text": mn.group(2)}
                )
                continue

            # --- blockquote ---
            if s.startswith("> "):
                cur["content"].append({"type": "quote", "text": s[2:]})
                continue

            # --- paragraph ---
            if s:
                cur["content"].append({"type": "paragraph", "text": s})

        if cur:
            items.append(cur)

    return items


def parse_notes(path):
    """Return {slide_number: text} from the script file."""
    if not path.exists():
        return {}
    notes = {}
    lines = path.read_text(encoding="utf-8").split("\n")
    cur = None
    buf = []
    for line in lines:
        m = re.match(r"^### Slide (\d+)", line)
        if m:
            if cur is not None and buf:
                notes[cur] = " ".join(buf)
            cur = int(m.group(1))
            buf = []
            continue
        if cur is not None:
            s = line.strip()
            if s.startswith(">") or s == "---" or s.startswith("###") or s.startswith("##"):
                continue
            if s:
                buf.append(s)
    if cur is not None and buf:
        notes[cur] = " ".join(buf)
    return notes


# ---------------------------------------------------------------------------
# Slide duplication
# ---------------------------------------------------------------------------

def dup(prs, tidx):
    """Duplicate template slide *tidx* preserving decorative shapes."""
    tmpl = prs.slides[tidx]
    new = prs.slides.add_slide(tmpl.slide_layout)

    old_c = tmpl._element.find(f"{{{NS_P}}}cSld")
    new_c = new._element.find(f"{{{NS_P}}}cSld")

    old_sp = old_c.find(f"{{{NS_P}}}spTree")
    new_sp = new_c.find(f"{{{NS_P}}}spTree")
    if old_sp is not None and new_sp is not None:
        for child in list(new_sp):
            new_sp.remove(child)
        for child in old_sp:
            new_sp.append(deepcopy(child))

    old_bg = old_c.find(f"{{{NS_P}}}bg")
    if old_bg is not None:
        new_bg = new_c.find(f"{{{NS_P}}}bg")
        if new_bg is not None:
            new_c.replace(new_bg, deepcopy(old_bg))
        else:
            new_c.insert(0, deepcopy(old_bg))

    for rel in tmpl.part.rels.values():
        if "image" in rel.reltype:
            try:
                new.part.relate_to(rel.target_part, rel.reltype, rel.rId)
            except Exception:
                pass

    return new


# ---------------------------------------------------------------------------
# Content rendering helpers
# ---------------------------------------------------------------------------

def _run(para, text, sz=13, bold=False, italic=False, color=DARK, font="Calibri"):
    r = para.add_run()
    r.text = text
    r.font.size = Pt(sz)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = font
    return r


def _rich(para, text, sz=13, base_bold=False, color=DARK):
    """Add rich-text segments to *para*."""
    for txt, b, i, c in parse_rich(text):
        _run(para, txt, sz=sz, bold=base_bold or b, italic=i, color=color,
             font="Consolas" if c else "Calibri")


def _txbox(slide, l, t, w, h):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tb.text_frame.word_wrap = True
    return tb


def add_title(slide, text, l, t, w, h, sz=22, color=NAVY, align=PP_ALIGN.LEFT):
    tb = _txbox(slide, l, t, w, h)
    p = tb.text_frame.paragraphs[0]
    p.alignment = align
    _run(p, strip_md(text), sz=sz, bold=True, color=color)


def add_number(slide, n):
    tb = _txbox(slide, 18.2, 10.5, 1.3, 0.5)
    p = tb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    _run(p, str(n), sz=10, color=GRAY)


def add_notes(slide, text):
    if text:
        slide.notes_slide.notes_text_frame.text = text


# ---------------------------------------------------------------------------
# Render content items into a rectangular zone
# ---------------------------------------------------------------------------

TEXT_TYPES = {"paragraph", "bullet", "numbered", "subheading", "heading1", "heading2", "quote"}


def render(slide, items, l, t, w, mh, sz=13):
    """Lay out *items* top-to-bottom inside (l, t, w, mh) zone."""
    y = t
    idx = 0

    while idx < len(items):
        remaining = mh - (y - t)
        if remaining < 0.25:
            break
        item = items[idx]

        # ---- text group ----
        if item["type"] in TEXT_TYPES:
            grp = []
            while idx < len(items) and items[idx]["type"] in TEXT_TYPES:
                grp.append(items[idx])
                idx += 1
            est = sum(1.8 if g["type"] in ("subheading", "heading1", "heading2") else 1.1 for g in grp)
            h = min(est * 0.28 + 0.2, remaining)
            tb = _txbox(slide, l, y, w, h)
            tf = tb.text_frame
            first = True
            for g in grp:
                p = tf.paragraphs[0] if first else tf.add_paragraph()
                first = False
                tp = g["type"]
                if tp in ("subheading", "heading2"):
                    p.space_before = Pt(8)
                    p.space_after = Pt(2)
                    _run(p, strip_md(g["text"]), sz=sz + 1, bold=True, color=NAVY)
                elif tp == "heading1":
                    p.alignment = PP_ALIGN.CENTER
                    p.space_before = Pt(4)
                    _run(p, strip_md(g["text"]), sz=22, bold=True, color=NAVY)
                elif tp == "bullet":
                    p.space_before = Pt(2)
                    p.space_after = Pt(2)
                    _rich(p, "• " + g["text"], sz=sz, color=DARK)
                elif tp == "numbered":
                    p.space_before = Pt(2)
                    p.space_after = Pt(2)
                    _rich(p, f"{g['num']}. {g['text']}", sz=sz, color=DARK)
                elif tp == "quote":
                    p.space_before = Pt(2)
                    _rich(p, g["text"], sz=sz - 1, color=GRAY)
                else:
                    p.space_before = Pt(2)
                    p.space_after = Pt(2)
                    _rich(p, g["text"], sz=sz, color=DARK)
            y += h + 0.1

        # ---- table ----
        elif item["type"] == "table":
            rows = item["rows"]
            if not rows:
                idx += 1
                continue
            nc = max(len(r) for r in rows)
            nr = len(rows)
            h = min(nr * 0.30 + 0.15, remaining, 4.0)
            tbl = slide.shapes.add_table(nr, nc, Inches(l), Inches(y), Inches(w), Inches(h)).table
            for ri, row in enumerate(rows):
                for ci in range(nc):
                    ct = strip_md(row[ci]) if ci < len(row) else ""
                    cell = tbl.cell(ri, ci)
                    cell.text = ct
                    for cp in cell.text_frame.paragraphs:
                        cp.font.size = Pt(10)
                        cp.font.color.rgb = DARK
                        if ri == 0:
                            cp.font.bold = True
                            cp.font.color.rgb = NAVY
            y += h + 0.12
            idx += 1

        # ---- code ----
        elif item["type"] == "code":
            lines = item["text"].split("\n")
            if len(lines) > MAX_CODE_LINES:
                lines = lines[:MAX_CODE_LINES] + ["// ..."]
            n = len(lines)
            h = min(n * 0.19 + 0.25, remaining, 6.0)
            tb = _txbox(slide, l, y, w, h)
            tb.fill.solid()
            tb.fill.fore_color.rgb = CODE_BG
            tf = tb.text_frame
            for li, ln in enumerate(lines):
                p = tf.paragraphs[0] if li == 0 else tf.add_paragraph()
                _run(p, ln, sz=9, color=DARK, font="Consolas")
                p.space_before = Pt(0)
                p.space_after = Pt(0)
            y += h + 0.1
            idx += 1

        # ---- image ----
        elif item["type"] == "image":
            h = min(remaining, 4.5)
            if h < 0.5:
                idx += 1
                continue
            try:
                from PIL import Image as PILImage

                img = PILImage.open(item["path"])
                iw, ih = img.size
                ratio = iw / ih
                pw, ph = Inches(w), Inches(h)
                if pw / ph > ratio:
                    act_h = ph
                    act_w = int(ph * ratio)
                else:
                    act_w = pw
                    act_h = int(pw / ratio)
                cx = Inches(l) + (pw - act_w) // 2
                slide.shapes.add_picture(item["path"], cx, Inches(y), act_w, act_h)
            except Exception as e:
                print(f"  Warning: image {item['path']}: {e}")
            y += h + 0.1
            idx += 1
        else:
            idx += 1


# ---------------------------------------------------------------------------
# Build
# ---------------------------------------------------------------------------

def choose_template(sd, is_last=False):
    if is_last or sd["number"] >= 49:
        return T_CONCL
    if sd["has_code"]:
        return random.choice(T_SPLIT)
    return random.choice(T_MAIN)


def build():
    add_sec = "--no-sections" not in sys.argv

    print("Parsing markdown…")
    items = parse_slides(SLIDES_DIR)
    notes = parse_notes(NOTES_FILE)
    slides = [i for i in items if i["type"] == "slide"]
    sections = [i for i in items if i["type"] == "section"]
    print(f"  {len(slides)} slides, {len(sections)} sections, {len(notes)} notes")

    print("Opening template…")
    prs = Presentation(str(TEMPLATE))
    n_tmpl = len(prs.slides)
    print(f"  {n_tmpl} template slides")

    sn = 0

    # ── title slide ──
    ts = dup(prs, T_TITLE)
    sn += 1
    add_title(ts, "HappyRow", 4.0, 2.5, 12.0, 2.0, sz=36, color=NAVY, align=PP_ALIGN.CENTER)
    tb = _txbox(ts, 4.0, 5.0, 12.0, 4.0)
    p0 = tb.text_frame.paragraphs[0]
    p0.alignment = PP_ALIGN.CENTER
    _run(p0, "Plateforme collaborative de gestion d'événements", sz=18, color=TEAL)
    p1 = tb.text_frame.add_paragraph()
    p1.alignment = PP_ALIGN.CENTER
    p1.space_before = Pt(16)
    _run(p1, "Titre Professionnel CDA — Niveau 6", sz=14, color=DARK)
    p2 = tb.text_frame.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(8)
    _run(p2, "[Nom Prénom] · [Date de session]", sz=12, color=GRAY)
    add_number(ts, sn)
    if 1 in notes:
        add_notes(ts, notes[1])

    # ── process all items ──
    for item in items:
        if item["type"] == "section":
            if add_sec and item["title"] != "Introduction":
                sn += 1
                sec = dup(prs, T_SECTION)
                add_title(sec, item["title"], 2.0, 3.5, 16.0, 4.0, sz=28, color=WHITE, align=PP_ALIGN.CENTER)
                add_number(sec, sn)
            continue

        if item["type"] != "slide" or item["number"] == 1:
            continue

        tidx = choose_template(item)
        ns_ = dup(prs, tidx)
        sn += 1

        is_split = tidx in T_SPLIT

        # title
        if tidx == T_SECTION:
            add_title(ns_, item["title"], 2.0, 3.5, 16.0, 4.0, sz=28, color=WHITE, align=PP_ALIGN.CENTER)
        else:
            add_title(ns_, item["title"], 1.0, 0.3, 12.0, 0.7, sz=22, color=NAVY)

        # content
        if is_split and item["has_code"]:
            code = [c for c in item["content"] if c["type"] == "code"]
            rest = [c for c in item["content"] if c["type"] != "code"]
            if tidx == T_SPLIT[0]:
                render(ns_, code, 0.8, 1.4, 8.8, 8.5, sz=12)
                render(ns_, rest, 10.2, 1.4, 8.8, 8.5, sz=12)
            else:
                render(ns_, rest, 0.8, 1.4, 8.8, 8.5, sz=12)
                render(ns_, code, 10.2, 1.4, 8.8, 8.5, sz=12)
        elif tidx == T_CONCL:
            render(ns_, item["content"], 1.0, 1.4, 17.5, 6.5)
        else:
            render(ns_, item["content"], 1.0, 1.4, 17.5, 8.5)

        note = notes.get(item["number"])
        if note:
            add_notes(ns_, note)

        add_number(ns_, sn)

    # ── remove original template slides ──
    print(f"Removing {n_tmpl} original template slides…")
    sldIdLst = prs.element.find(f"{{{NS_P}}}sldIdLst")
    originals = list(sldIdLst)[:n_tmpl]
    for sid in originals:
        rid = sid.get(f"{{{NS_R}}}id")
        sldIdLst.remove(sid)
        if rid:
            try:
                prs.part.drop_rel(rid)
            except Exception:
                pass

    prs.save(str(OUTPUT))
    total = len(prs.slides)
    notes_count = sum(
        1 for sl in prs.slides
        if sl.has_notes_slide and sl.notes_slide.notes_text_frame.text.strip()
    )
    print(f"\nDone! {total} slides ({notes_count} with notes) → {OUTPUT.name}")


if __name__ == "__main__":
    build()
