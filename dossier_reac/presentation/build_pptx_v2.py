#!/usr/bin/env python3
"""Build HappyRow-Presentation.pptx using the repaired native-layout template.

Usage: python3 build_pptx_v2.py [--no-sections]

Reads markdown from slides/, oral script from 00_SCRIPT_COMPLET.md,
and generates a 50-slide presentation with speaker notes.
"""

import os
import re
import random
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

SCRIPT_DIR = Path(__file__).parent
SLIDES_DIR = SCRIPT_DIR / "slides"
TEMPLATE = SCRIPT_DIR / "ppt" / "HappyRow-Template  -  Repaired.pptx"
NOTES_FILE = SCRIPT_DIR / "00_SCRIPT_COMPLET.md"
IMG_DIR = SCRIPT_DIR / "mermaid-images"
OUTPUT = SCRIPT_DIR / "HappyRow-Presentation.pptx"

NAVY = RGBColor(0x3D, 0x5A, 0x6C)
TEAL = RGBColor(0x5F, 0xBD, 0xB4)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x2D, 0x3A, 0x3A)
GRAY = RGBColor(0x99, 0x99, 0x99)
CODE_BG = RGBColor(0xF5, 0xF5, 0xF5)

MAX_CODE_LINES = 22

# Layout indices in the repaired template
L_MAIN = [0, 1, 2]       # MAIN 1/2/3  — Title + Content
L_SPLIT = [3, 4]          # SPLIT 1/2   — Title + Left + Right
L_SPLIT13 = 5             # SPLIT 1-3   — Title + Left(1/3) + Right(2/3)
L_SECTION = 6             # SECTION     — Center Title
L_COVER = 7               # COVER       — Center Title + Subtitle

# Slide dimensions (EMU)
SLD_W_IN = 20.32  # inches (18288000 EMU)
SLD_H_IN = 11.43  # inches (10287000 EMU)

random.seed(42)


# ---------------------------------------------------------------------------
# Markdown parsing (reused from build_pptx.py)
# ---------------------------------------------------------------------------

def parse_rich(text):
    segs = []
    pat = re.compile(r"(\*\*(.+?)\*\*)|(\*(.+?)\*)|(`(.+?)`)")
    pos = 0
    for m in pat.finditer(text):
        if m.start() > pos:
            segs.append((text[pos:m.start()], False, False, False))
        if m.group(2):
            segs.append((m.group(2), True, False, False))
        elif m.group(4):
            segs.append((m.group(4), False, True, False))
        elif m.group(6):
            segs.append((m.group(6), False, False, True))
        pos = m.end()
    if pos < len(text):
        segs.append((text[pos:], False, False, False))
    return segs or [(text, False, False, False)]


def strip_md(text):
    t = re.sub(r"\*\*(.+?)\*\*", r"\1", text)
    t = re.sub(r"\*(.+?)\*", r"\1", t)
    t = re.sub(r"`(.+?)`", r"\1", t)
    return t.strip()


def parse_slides(slides_dir):
    items = []
    mermaid_idx = 0

    for md in sorted(Path(slides_dir).glob("*.md")):
        lines = md.read_text(encoding="utf-8").split("\n")
        cur = None
        in_code = in_mermaid = False
        code_lang = ""
        code_buf = []
        file_section_seen = False

        for raw in lines:
            s = raw.strip()
            if s.startswith("<!--") or s == "---":
                continue
            if s.startswith("> *[") or s.startswith(">*["):
                continue

            mh = re.match(r"^# (.+)$", s)
            if mh and not s.startswith("## ") and not in_code and not in_mermaid:
                if not file_section_seen:
                    file_section_seen = True
                    if cur:
                        items.append(cur)
                        cur = None
                    items.append({"type": "section", "title": mh.group(1)})
                elif cur:
                    cur["content"].append({"type": "heading1", "text": mh.group(1)})
                continue

            ms = re.match(r"^## Slide (\d+)\s*[—–-]\s*(.+?)(?:\s*\(.+\))?$", s)
            if ms and not in_code and not in_mermaid:
                if cur:
                    items.append(cur)
                cur = {
                    "type": "slide",
                    "number": int(ms.group(1)),
                    "title": ms.group(2).strip(),
                    "content": [],
                    "has_code": False,
                    "has_image": False,
                }
                continue

            if cur is None:
                continue

            if s == "```mermaid" and not in_code:
                in_mermaid = True
                mermaid_idx += 1
                continue
            if in_mermaid:
                if s == "```":
                    in_mermaid = False
                    p = str(IMG_DIR / f"diagram_{mermaid_idx}.png")
                    if os.path.exists(p):
                        cur["content"].append({"type": "image", "path": p})
                        cur["has_image"] = True
                continue

            mc = re.match(r"^```(\w+)", s)
            if mc and not in_code:
                in_code = True
                code_lang = mc.group(1)
                code_buf = []
                continue
            if s == "```" and in_code:
                in_code = False
                cur["has_code"] = True
                cur["content"].append(
                    {"type": "code", "language": code_lang, "text": "\n".join(code_buf)}
                )
                continue
            if in_code:
                code_buf.append(raw)
                continue

            if s == "```":
                continue

            if s.startswith("|") and s.endswith("|"):
                if re.match(r"^\|[\s\-:|]+\|$", s):
                    continue
                cells = [c.strip() for c in s.split("|")[1:-1]]
                if cur["content"] and cur["content"][-1]["type"] == "table":
                    cur["content"][-1]["rows"].append(cells)
                else:
                    cur["content"].append({"type": "table", "rows": [cells]})
                continue

            m3 = re.match(r"^### (.+)$", s)
            if m3:
                cur["content"].append({"type": "subheading", "text": m3.group(1)})
                continue

            m2 = re.match(r"^## (.+)$", s)
            if m2:
                cur["content"].append({"type": "heading2", "text": m2.group(1)})
                continue

            if s.startswith("- "):
                cur["content"].append({"type": "bullet", "text": s[2:]})
                continue

            mn = re.match(r"^(\d+)\.\s+(.+)$", s)
            if mn:
                cur["content"].append(
                    {"type": "numbered", "num": mn.group(1), "text": mn.group(2)}
                )
                continue

            if s.startswith("> "):
                cur["content"].append({"type": "quote", "text": s[2:]})
                continue

            if s:
                cur["content"].append({"type": "paragraph", "text": s})

        if cur:
            items.append(cur)

    return items


def parse_notes(path):
    if not path.exists():
        return {}
    notes = {}
    lines = path.read_text(encoding="utf-8").split("\n")
    cur = None
    buf = []
    for line in lines:
        m = re.match(r"^### Slide (\d+)", line)
        if m:
            if cur is not None and buf:
                notes[cur] = "\n".join(buf)
            cur = int(m.group(1))
            buf = []
            continue
        if cur is not None:
            s = line.strip()
            if s.startswith("###") or s.startswith("##"):
                continue
            if s.startswith("> "):
                buf.append(f"[{s[2:]}]")
                continue
            if s == "---":
                continue
            if s:
                buf.append(s)
    if cur is not None and buf:
        notes[cur] = "\n".join(buf)
    return notes


# ---------------------------------------------------------------------------
# Rendering helpers
# ---------------------------------------------------------------------------

def _run(para, text, sz=13, bold=False, italic=False, color=DARK, font="Calibri"):
    r = para.add_run()
    r.text = text
    r.font.size = Pt(sz)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = font
    return r


def _rich(para, text, sz=13, base_bold=False, color=DARK):
    for txt, b, i, c in parse_rich(text):
        _run(para, txt, sz=sz, bold=base_bold or b, italic=i, color=color,
             font="Consolas" if c else "Calibri")


def _txbox(slide, l, t, w, h):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tb.text_frame.word_wrap = True
    return tb


TEXT_TYPES = {"paragraph", "bullet", "numbered", "subheading", "heading1", "heading2", "quote"}


def render_into_textframe(tf, items, sz=13):
    """Render text-type items into an existing text frame (placeholder)."""
    first = True
    for item in items:
        tp = item["type"]
        if tp not in TEXT_TYPES:
            continue
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        if tp in ("subheading", "heading2"):
            p.space_before = Pt(8)
            p.space_after = Pt(2)
            _run(p, strip_md(item["text"]), sz=sz + 1, bold=True, color=NAVY)
        elif tp == "heading1":
            p.alignment = PP_ALIGN.CENTER
            p.space_before = Pt(4)
            _run(p, strip_md(item["text"]), sz=20, bold=True, color=NAVY)
        elif tp == "bullet":
            p.space_before = Pt(2)
            p.space_after = Pt(2)
            _rich(p, "• " + item["text"], sz=sz, color=DARK)
        elif tp == "numbered":
            p.space_before = Pt(2)
            p.space_after = Pt(2)
            _rich(p, f"{item['num']}. {item['text']}", sz=sz, color=DARK)
        elif tp == "quote":
            p.space_before = Pt(2)
            _rich(p, item["text"], sz=sz - 1, color=GRAY)
        else:
            p.space_before = Pt(2)
            p.space_after = Pt(2)
            _rich(p, item["text"], sz=sz, color=DARK)


def render_code_textframe(tf, code_text, sz=9):
    """Render code block into a text frame."""
    lines = code_text.split("\n")
    if len(lines) > MAX_CODE_LINES:
        lines = lines[:MAX_CODE_LINES] + ["// ..."]
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        _run(p, ln, sz=sz, color=DARK, font="Consolas")
        p.space_before = Pt(0)
        p.space_after = Pt(0)


def render_table(slide, rows, l, t, w, max_h):
    """Add a table shape to the slide."""
    if not rows:
        return
    nc = max(len(r) for r in rows)
    nr = len(rows)
    h = min(nr * 0.30 + 0.15, max_h, 4.0)
    tbl = slide.shapes.add_table(nr, nc, Inches(l), Inches(t), Inches(w), Inches(h)).table
    for ri, row in enumerate(rows):
        for ci in range(nc):
            ct = strip_md(row[ci]) if ci < len(row) else ""
            cell = tbl.cell(ri, ci)
            cell.text = ct
            for cp in cell.text_frame.paragraphs:
                cp.font.size = Pt(10)
                cp.font.color.rgb = DARK
                if ri == 0:
                    cp.font.bold = True
                    cp.font.color.rgb = NAVY


def render_image(slide, img_path, l, t, w, max_h):
    """Add an image to the slide, preserving aspect ratio."""
    try:
        from PIL import Image as PILImage
        img = PILImage.open(img_path)
        iw, ih = img.size
        ratio = iw / ih
        pw, ph = Inches(w), Inches(max_h)
        if pw / ph > ratio:
            act_h = ph
            act_w = int(ph * ratio)
        else:
            act_w = pw
            act_h = int(pw / ratio)
        cx = Inches(l) + (pw - act_w) // 2
        slide.shapes.add_picture(img_path, cx, Inches(t), act_w, act_h)
    except Exception as e:
        print(f"  Warning: image {img_path}: {e}")


def render_extras(slide, items, l, t, w, max_h, sz=12):
    """Render non-text items (tables, code, images) as shapes on the slide."""
    y = t
    for item in items:
        remaining = max_h - (y - t)
        if remaining < 0.3:
            break
        if item["type"] == "table":
            h = min(len(item["rows"]) * 0.30 + 0.15, remaining, 4.0)
            render_table(slide, item["rows"], l, y, w, remaining)
            y += h + 0.12
        elif item["type"] == "code":
            lines = item["text"].split("\n")
            if len(lines) > MAX_CODE_LINES:
                lines = lines[:MAX_CODE_LINES] + ["// ..."]
            n = len(lines)
            h = min(n * 0.19 + 0.25, remaining, 6.0)
            tb = _txbox(slide, l, y, w, h)
            tb.fill.solid()
            tb.fill.fore_color.rgb = CODE_BG
            render_code_textframe(tb.text_frame, item["text"], sz=9)
            y += h + 0.1
        elif item["type"] == "image":
            h = min(remaining, 4.5)
            if h >= 0.5:
                render_image(slide, item["path"], l, y, w, h)
                y += h + 0.1


def render_full(slide, items, l, t, w, max_h, sz=13):
    """Render ALL items (text + extras) sequentially as shapes."""
    y = t
    idx = 0
    while idx < len(items):
        remaining = max_h - (y - t)
        if remaining < 0.25:
            break
        item = items[idx]

        if item["type"] in TEXT_TYPES:
            grp = []
            while idx < len(items) and items[idx]["type"] in TEXT_TYPES:
                grp.append(items[idx])
                idx += 1
            est = sum(1.8 if g["type"] in ("subheading", "heading1", "heading2") else 1.1 for g in grp)
            h = min(est * 0.28 + 0.2, remaining)
            tb = _txbox(slide, l, y, w, h)
            render_into_textframe(tb.text_frame, grp, sz=sz)
            y += h + 0.1
        elif item["type"] == "table":
            h = min(len(item["rows"]) * 0.30 + 0.15, remaining, 4.0)
            render_table(slide, item["rows"], l, y, w, remaining)
            y += h + 0.12
            idx += 1
        elif item["type"] == "code":
            lines = item["text"].split("\n")
            if len(lines) > MAX_CODE_LINES:
                lines = lines[:MAX_CODE_LINES] + ["// ..."]
            n = len(lines)
            h = min(n * 0.19 + 0.25, remaining, 6.0)
            tb = _txbox(slide, l, y, w, h)
            tb.fill.solid()
            tb.fill.fore_color.rgb = CODE_BG
            render_code_textframe(tb.text_frame, item["text"], sz=9)
            y += h + 0.1
            idx += 1
        elif item["type"] == "image":
            h = min(remaining, 4.5)
            if h >= 0.5:
                render_image(slide, item["path"], l, y, w, h)
            y += h + 0.1
            idx += 1
        else:
            idx += 1


# ---------------------------------------------------------------------------
# Layout selection
# ---------------------------------------------------------------------------

_main_cycle = 0
_split_cycle = 0


def choose_layout(sd, is_last=False):
    global _main_cycle, _split_cycle
    if sd["number"] == 1:
        return L_COVER
    if is_last or sd["number"] >= 49:
        li = L_MAIN[_main_cycle % len(L_MAIN)]
        _main_cycle += 1
        return li
    if sd["has_code"]:
        li = L_SPLIT[_split_cycle % len(L_SPLIT)]
        _split_cycle += 1
        return li
    if sd["has_image"]:
        return L_SPLIT13
    li = L_MAIN[_main_cycle % len(L_MAIN)]
    _main_cycle += 1
    return li


# ---------------------------------------------------------------------------
# Build
# ---------------------------------------------------------------------------

NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"
NS_R = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"


def remove_existing_slides(prs):
    """Remove all slides that came with the template."""
    sld_id_lst = prs.element.find(f"{{{NS_P}}}sldIdLst")
    if sld_id_lst is None:
        return
    for sid in list(sld_id_lst):
        rid = sid.get(f"{{{NS_R}}}id")
        sld_id_lst.remove(sid)
        if rid:
            try:
                prs.part.drop_rel(rid)
            except Exception:
                pass


def build():
    add_sec = "--no-sections" not in sys.argv

    print("Parsing markdown…")
    items = parse_slides(SLIDES_DIR)
    notes = parse_notes(NOTES_FILE)
    slides = [i for i in items if i["type"] == "slide"]
    sections = [i for i in items if i["type"] == "section"]
    print(f"  {len(slides)} slides, {len(sections)} sections, {len(notes)} notes")

    print(f"Opening template: {TEMPLATE.name}")
    prs = Presentation(str(TEMPLATE))
    remove_existing_slides(prs)
    print(f"  {len(prs.slide_layouts)} layouts available")

    sn = 0

    for item in items:

        # ── Section dividers ──
        if item["type"] == "section":
            if add_sec and item["title"] != "Introduction":
                sn += 1
                layout = prs.slide_layouts[L_SECTION]
                slide = prs.slides.add_slide(layout)
                ph = slide.placeholders[0]
                ph.text = ""
                p = ph.text_frame.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER
                _run(p, strip_md(item["title"]), sz=28, bold=True, color=WHITE)
                print(f"  [{sn:2d}] SECTION: {item['title']}")
            continue

        if item["type"] != "slide":
            continue

        sd = item
        is_last = sd["number"] >= 49
        li = choose_layout(sd, is_last)
        layout = prs.slide_layouts[li]
        slide = prs.slides.add_slide(layout)
        sn += 1

        layout_name = layout.name

        # ── COVER slide ──
        if li == L_COVER:
            ph_title = slide.placeholders[0]
            ph_title.text = ""
            p = ph_title.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            _run(p, "HappyRow", sz=36, bold=True, color=NAVY)

            ph_sub = slide.placeholders[10]
            ph_sub.text = ""
            tf = ph_sub.text_frame
            p0 = tf.paragraphs[0]
            p0.alignment = PP_ALIGN.CENTER
            _run(p0, "Plateforme collaborative de gestion d'événements", sz=18, color=TEAL)
            p1 = tf.add_paragraph()
            p1.alignment = PP_ALIGN.CENTER
            p1.space_before = Pt(16)
            _run(p1, "Titre Professionnel CDA — Niveau 6", sz=14, color=DARK)
            p2 = tf.add_paragraph()
            p2.alignment = PP_ALIGN.CENTER
            p2.space_before = Pt(8)
            _run(p2, "Jimmy Ni · 24/03/2026 · Simplon", sz=12, color=GRAY)

            print(f"  [{sn:2d}] COVER")

        # ── SPLIT slides (code on one side, text on the other) ──
        elif li in L_SPLIT:
            ph_title = slide.placeholders[0]
            ph_title.text = ""
            _run(ph_title.text_frame.paragraphs[0], strip_md(sd["title"]),
                 sz=20, bold=True, color=NAVY)

            code_items = [c for c in sd["content"] if c["type"] == "code"]
            text_items = [c for c in sd["content"] if c["type"] in TEXT_TYPES]
            other_items = [c for c in sd["content"]
                           if c["type"] not in ("code",) and c["type"] not in TEXT_TYPES]

            # Left placeholder (idx=10): code
            ph_left = slide.placeholders[10]
            ph_left.text = ""
            if code_items:
                render_code_textframe(ph_left.text_frame, code_items[0]["text"], sz=8)
                ph_left.text_frame.paragraphs[0].font.size = Pt(8)

            # Right placeholder (idx=11): text
            ph_right = slide.placeholders[11]
            ph_right.text = ""
            if text_items:
                render_into_textframe(ph_right.text_frame, text_items, sz=11)

            # Extra tables/images below
            if other_items:
                render_extras(slide, other_items, 1.0, 8.5, 18.0, 2.5, sz=10)

            print(f"  [{sn:2d}] {layout_name}: {sd['title'][:50]}")

        # ── SPLIT 1-3 (image/diagram on left, text on right) ──
        elif li == L_SPLIT13:
            ph_title = slide.placeholders[0]
            ph_title.text = ""
            _run(ph_title.text_frame.paragraphs[0], strip_md(sd["title"]),
                 sz=20, bold=True, color=NAVY)

            img_items = [c for c in sd["content"] if c["type"] == "image"]
            text_items = [c for c in sd["content"] if c["type"] in TEXT_TYPES]
            table_items = [c for c in sd["content"] if c["type"] == "table"]
            code_items = [c for c in sd["content"] if c["type"] == "code"]

            # Left (1/3): image or first special item
            ph_left = slide.placeholders[10]
            ph_left.text = ""
            if img_items:
                left = ph_left.left
                top = ph_left.top
                w_emu = ph_left.width
                h_emu = ph_left.height
                render_image(slide, img_items[0]["path"],
                             left / 914400, top / 914400,
                             w_emu / 914400, h_emu / 914400)
            elif text_items:
                half = len(text_items) // 2
                render_into_textframe(ph_left.text_frame, text_items[:max(half, 1)], sz=11)

            # Right (2/3): text
            ph_right = slide.placeholders[11]
            ph_right.text = ""
            remaining_text = text_items if img_items else text_items[max(len(text_items)//2, 1):]
            if remaining_text:
                render_into_textframe(ph_right.text_frame, remaining_text, sz=11)

            # Extra tables/code
            extras = table_items + code_items
            if extras:
                render_extras(slide, extras, 1.0, 8.5, 18.0, 2.5, sz=10)

            print(f"  [{sn:2d}] {layout_name}: {sd['title'][:50]}")

        # ── MAIN slides (title + full content) ──
        else:
            ph_title = slide.placeholders[0]
            ph_title.text = ""
            _run(ph_title.text_frame.paragraphs[0], strip_md(sd["title"]),
                 sz=20, bold=True, color=NAVY)

            text_items = [c for c in sd["content"] if c["type"] in TEXT_TYPES]
            extra_items = [c for c in sd["content"] if c["type"] not in TEXT_TYPES]

            ph_content = slide.placeholders[10]
            ph_content.text = ""
            if text_items:
                render_into_textframe(ph_content.text_frame, text_items, sz=12)

            if extra_items:
                # Place extras below the content placeholder
                ph_bot = ph_content.top + ph_content.height
                bot_y = ph_bot / 914400
                avail = SLD_H_IN - bot_y - 0.3
                if avail > 0.3:
                    render_extras(slide, extra_items, 1.0, bot_y, 18.0, avail, sz=10)
                else:
                    render_extras(slide, extra_items, 1.0, 2.0, 18.0, 8.0, sz=10)

            print(f"  [{sn:2d}] {layout_name}: {sd['title'][:50]}")

        # ── Speaker notes ──
        note = notes.get(sd["number"])
        if note:
            slide.notes_slide.notes_text_frame.text = note

    # ── Summary ──
    total = len(prs.slides)
    notes_count = sum(
        1 for sl in prs.slides
        if sl.has_notes_slide and sl.notes_slide.notes_text_frame.text.strip()
    )
    print(f"\nSaving: {OUTPUT.name}")
    prs.save(str(OUTPUT))
    print(f"Done! {total} slides ({notes_count} with speaker notes)")


if __name__ == "__main__":
    build()
