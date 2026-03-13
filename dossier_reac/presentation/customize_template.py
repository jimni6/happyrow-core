#!/usr/bin/env python3
"""Customize the pandoc reference.pptx template with HappyRow brand colors."""

from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import copy
from lxml import etree

NAVY = RGBColor(0x3D, 0x5A, 0x6C)
TEAL = RGBColor(0x5F, 0xBD, 0xB4)
CORAL = RGBColor(0xE6, 0xA1, 0x9A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x2D, 0x3A, 0x3A)
LIGHT_BG = RGBColor(0xF7, 0xF9, 0xFA)

def customize_theme_colors(prs):
    """Modify the theme XML to set HappyRow colors."""
    ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
    theme = prs.slide_masters[0].slide_layouts[0].slide_master.element
    theme_elements = theme.findall('.//' + '{http://schemas.openxmlformats.org/drawingml/2006/main}clrScheme')
    
    for clr_scheme in theme_elements:
        for dk1 in clr_scheme.findall('a:dk1/a:srgbClr', ns):
            dk1.set('val', '2D3A3A')
        for dk2 in clr_scheme.findall('a:dk2/a:srgbClr', ns):
            dk2.set('val', '3D5A6C')
        for lt1 in clr_scheme.findall('a:lt1/a:srgbClr', ns):
            lt1.set('val', 'FFFFFF')
        for lt2 in clr_scheme.findall('a:lt2/a:srgbClr', ns):
            lt2.set('val', 'F7F9FA')
        for accent1 in clr_scheme.findall('a:accent1/a:srgbClr', ns):
            accent1.set('val', '5FBDB4')
        for accent2 in clr_scheme.findall('a:accent2/a:srgbClr', ns):
            accent2.set('val', 'E6A19A')
        for accent3 in clr_scheme.findall('a:accent3/a:srgbClr', ns):
            accent3.set('val', '3D5A6C')
        for accent4 in clr_scheme.findall('a:accent4/a:srgbClr', ns):
            accent4.set('val', '5FBDB4')

def main():
    prs = Presentation('reference.pptx')

    try:
        customize_theme_colors(prs)
    except Exception as e:
        print(f"Note: Could not modify theme XML directly ({e}), continuing with layout styling...")

    for master in prs.slide_masters:
        for layout in master.slide_layouts:
            for ph in layout.placeholders:
                if ph.placeholder_format.idx == 0:
                    for para in ph.text_frame.paragraphs:
                        for run in para.runs:
                            run.font.color.rgb = NAVY
                            run.font.size = Pt(32)
                elif ph.placeholder_format.idx == 1:
                    for para in ph.text_frame.paragraphs:
                        for run in para.runs:
                            run.font.color.rgb = DARK
                            run.font.size = Pt(18)

    prs.save('reference.pptx')
    print("Template customized with HappyRow colors (navy=#3D5A6C, teal=#5FBDB4, coral=#E6A19A)")

if __name__ == '__main__':
    main()
